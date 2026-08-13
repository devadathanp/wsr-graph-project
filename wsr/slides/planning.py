"""Slide 11 — Quarterly planning chart."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from wsr.slides.base import new_content_slide
from wsr_style import style_body_run


def add_planning_slide(
    prs: Presentation,
    report_date: str,
    qp: dict[str, int | str] | None,
    *,
    chart_image: Path | None = None,
    slide_number: int = 10,
    quarter_label: str = "Q3-2026",
    fiscal_year: int = 2026,
) -> None:
    slide = new_content_slide(
        prs,
        f"Quarterly Planning {fiscal_year}-Non STLA",
        report_date,
        slide_number,
    )

    if qp is None:
        note = slide.shapes.add_textbox(Inches(0.59), Inches(2.5), Inches(12.18), Inches(0.8))
        note_run = note.text_frame.paragraphs[0].add_run()
        note_run.text = "No Actual Available Estimate found on the Scrum sheet."
        style_body_run(note_run)
        return

    if chart_image is not None and chart_image.exists():
        slide.shapes.add_picture(
            str(chart_image),
            Inches(1.29),
            Inches(0.93),
            width=Inches(11.20),
            height=Inches(4.91),
        )

    note = slide.shapes.add_textbox(Inches(0.59), Inches(6.05), Inches(12.18), Inches(0.40))
    note_run = note.text_frame.paragraphs[0].add_run()
    note_run.text = (
        f"Note: The Estimations for {quarter_label} Planning are in progress, "
        "these are high level tentative estimations"
    )
    style_body_run(note_run)
