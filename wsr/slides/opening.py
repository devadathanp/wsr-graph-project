"""Opening slides: title, agenda, and MOM."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from wsr.constants import AGENDA_BADGE_SIZE, AGENDA_ITEMS, AGENDA_LAYOUT, LAYOUT_OPENING
from wsr.slides.base import add_table, empty_row, new_content_slide
from wsr_style import (
    FONT_BODY,
    FONT_MAJOR,
    TEXT_DARK,
    TITLE_SIZE,
    add_number_badge,
    find_placeholder,
    set_run_font,
    set_slide_footer,
    style_agenda_run,
    style_title_date_run,
    style_title_run,
)


def add_title_slide(prs: Presentation, report_date: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_OPENING])

    title_ph = find_placeholder(slide, idx=0)
    if title_ph is not None:
        text_frame = title_ph.text_frame
        text_frame.clear()
        line1 = text_frame.paragraphs[0].add_run()
        line1.text = "CES PFS CSAR (Non- STLA) and CORE 2"
        style_title_run(line1)
        line2 = text_frame.add_paragraph().add_run()
        line2.text = "Weekly Status Report"
        set_run_font(line2, size=TITLE_SIZE, bold=False, color=TEXT_DARK, name=FONT_BODY)

    date_ph = find_placeholder(slide, idx=16)
    if date_ph is not None:
        date_ph.text = report_date
        for paragraph in date_ph.text_frame.paragraphs:
            for run in paragraph.runs:
                style_title_date_run(run)

    presenter_ph = find_placeholder(slide, idx=23)
    if presenter_ph is not None:
        presenter_ph.text = ""

    set_slide_footer(slide, report_date, 1)


def add_agenda_slide(prs: Presentation, report_date: str) -> None:
    slide = new_content_slide(prs, "Agenda", report_date, 2)
    for idx, (item, layout) in enumerate(zip(AGENDA_ITEMS, AGENDA_LAYOUT), start=1):
        add_number_badge(
            slide,
            str(idx),
            left=0.67,
            top=layout["badge_top"],
            accent=True,
            size=AGENDA_BADGE_SIZE,
        )
        box = slide.shapes.add_textbox(
            Inches(layout["text_left"]),
            Inches(layout["text_top"]),
            Inches(4.3),
            Inches(0.5),
        )
        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].add_run()
        run.text = item
        style_agenda_run(run)


def add_mom_slide(prs: Presentation, report_date: str) -> None:
    slide = new_content_slide(prs, "MOM of 11/06,18/06", report_date, 3)
    headers = [
        "Sr. No.",
        "Discussion points",
        "Action Items",
        "Status",
        "Ownership",
        "Action item identified",
        "Action closure date",
        "Remarks",
    ]
    rows = [empty_row(len(headers))]
    add_table(
        slide,
        headers,
        rows,
        col_widths=[0.55, 2.2, 1.45, 0.9, 0.95, 1.45, 1.35, 2.45],
    )
