"""Shared slide helpers: content slides and tables."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from wsr.constants import LAYOUT_CONTENT
from wsr_style import (
    TABLE_HEADER_ROW_HEIGHT_IN,
    TABLE_LEFT_IN,
    TABLE_WIDTH_IN,
    content_top_below_title,
    fit_table_column_widths,
    set_slide_footer,
    set_slide_title,
    style_key_value_table,
    style_table_cells,
)


def empty_row(column_count: int) -> list[str]:
    return ["-"] * column_count


def new_content_slide(
    prs: Presentation,
    title: str,
    report_date: str,
    slide_number: int,
    *,
    title_size=None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    if title_size is not None:
        set_slide_title(slide, title, size=title_size)
    else:
        set_slide_title(slide, title)
    set_slide_footer(slide, report_date, slide_number)
    return slide


def add_table(slide, headers, rows, top=None, col_widths=None) -> float:
    if top is None:
        top = content_top_below_title(slide)
    fitted_widths = fit_table_column_widths(headers, col_widths)
    row_height = min(0.34 * (len(rows) + 1), 5.8)
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(TABLE_LEFT_IN),
        Inches(top),
        Inches(TABLE_WIDTH_IN),
        Inches(row_height),
    )
    table = table_shape.table
    for idx, width in enumerate(fitted_widths):
        table.columns[idx].width = Inches(width)
    table.rows[0].height = Inches(TABLE_HEADER_ROW_HEIGHT_IN)

    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)

    style_table_cells(table)
    return top


def add_summary_key_value_table(
    slide,
    rows: list[tuple[str, str]],
    *,
    left: float,
    top: float,
    width: float,
) -> float:
    row_height = 0.32
    table_height = row_height * max(len(rows), 1)
    table_shape = slide.shapes.add_table(
        len(rows),
        2,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(table_height),
    )
    table = table_shape.table
    table.columns[0].width = Inches(width * 0.56)
    table.columns[1].width = Inches(width * 0.44)
    for row_idx, (label, value) in enumerate(rows):
        table.cell(row_idx, 0).text = label
        table.cell(row_idx, 1).text = str(value)
    style_key_value_table(table)
    return top + table_height
