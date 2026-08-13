"""Slide 10 — Risks & Mitigation Plan (from Risk and Mitigation Plan sheet)."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from wsr.slides.base import add_table, empty_row, new_content_slide
from wsr_style import TABLE_LEFT_IN, style_body_run


def _add_impact_legend(slide) -> None:
    legend_items = [
        (RGBColor(0xFF, 0x00, 0x00), "High Impact / High Possibility"),
        (RGBColor(0xFF, 0xBF, 0x00), "Medium Impact / Medium Possibility"),
        (RGBColor(0x00, 0xB0, 0x50), "Low Impact / Low Possibility"),
    ]
    circle_top = 6.36
    text_top = 6.33
    item_width = 3.95
    gap = 0.35
    start_left = TABLE_LEFT_IN
    circle_size = 0.24

    for index, (color, label) in enumerate(legend_items):
        left = start_left + index * (item_width + gap)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left),
            Inches(circle_top),
            Inches(circle_size),
            Inches(circle_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        box = slide.shapes.add_textbox(
            Inches(left + 0.34),
            Inches(text_top),
            Inches(item_width - 0.34),
            Inches(0.34),
        )
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].add_run()
        run.text = label
        style_body_run(run)


def add_risks_slide(
    prs: Presentation,
    report_date: str,
    items: list[dict] | None = None,
    *,
    slide_number: int = 9,
) -> None:
    slide = new_content_slide(prs, "Risks & Mitigation Plan", report_date, slide_number)
    headers = ["#", "DCR", "Risk / Issue", "Impact", "Support Required", "Status", "RAG Status"]
    rows = [
        [
            str(i + 1),
            item.get("dcr", ""),
            item.get("risk", ""),
            item.get("impact", ""),
            item.get("support", ""),
            item.get("status", ""),
            "",
        ]
        for i, item in enumerate(items or [])
    ]
    if not rows:
        rows = [empty_row(len(headers))]
        rows[0][-1] = ""

    add_table(
        slide,
        headers,
        rows,
        col_widths=[0.45, 0.95, 3.7, 1.45, 2.2, 1.15, 1.2],
    )
    _add_impact_legend(slide)
