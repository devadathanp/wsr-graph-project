"""Slide 4 — DCR status charts and summary panel."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from wsr.constants import (
    DCR_CHART_HEIGHT,
    DCR_CHART_LEFT,
    DCR_CHART_WIDTH,
    DCR_EVAL_TOP,
    DCR_IMPL_TOP,
    DCR_PANEL_LEFT,
    DCR_PANEL_WIDTH,
    DCR_SUMMARY_TOP,
    DCR_TITLE_TOP,
)
from wsr.slides.base import add_summary_key_value_table, new_content_slide
from wsr.tracker import format_ordinal_day_month, format_quarter_label
from wsr_style import raise_slide_title


def dcr_status_slide_title(report_date: str) -> str:
    quarter = format_quarter_label(report_date)
    till = format_ordinal_day_month(report_date)
    return (
        f"DCR Status {quarter} - CSAR (Non-STLA) & Core 2 program - PFS (till {till})"
    )


def add_dcr_status_slide(
    prs: Presentation,
    report_date: str,
    impl_chart: Path,
    eval_chart: Path,
    summary_rows: list[tuple[str, str]],
) -> None:
    slide = new_content_slide(
        prs,
        dcr_status_slide_title(report_date),
        report_date,
        4,
        title_size=Pt(18),
    )
    raise_slide_title(slide, top_in=DCR_TITLE_TOP)

    slide.shapes.add_picture(
        str(eval_chart),
        Inches(DCR_CHART_LEFT),
        Inches(DCR_EVAL_TOP),
        width=Inches(DCR_CHART_WIDTH),
        height=Inches(DCR_CHART_HEIGHT),
    )
    slide.shapes.add_picture(
        str(impl_chart),
        Inches(DCR_CHART_LEFT),
        Inches(DCR_IMPL_TOP),
        width=Inches(DCR_CHART_WIDTH),
        height=Inches(DCR_CHART_HEIGHT),
    )

    add_summary_key_value_table(
        slide,
        summary_rows,
        left=DCR_PANEL_LEFT,
        top=DCR_SUMMARY_TOP,
        width=DCR_PANEL_WIDTH,
    )
