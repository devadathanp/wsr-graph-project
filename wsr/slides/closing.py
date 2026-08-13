"""Closing slide — tree backdrop plus two lines of copy. Nothing else."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from wsr.constants import LAYOUT_BLANK
from wsr_style import (
    CLOSING_HEADLINE_SIZE,
    CLOSING_SUBLINE_SIZE,
    DEFAULT_CLOSING_BACKDROP,
    FONT_BODY,
    FONT_MAJOR,
    WHITE,
    set_run_font,
)


def _send_shape_to_back(shape) -> None:
    sp_tree = shape._element.getparent()
    if sp_tree is None:
        return
    sp_tree.remove(shape._element)
    insert_at = 1 if len(sp_tree) > 0 else 0
    sp_tree.insert(insert_at, shape._element)


def resolve_closing_backdrop(assets_dir: Path, override: Path | None = None) -> Path | None:
    if override is not None and override.exists():
        return override
    candidates = [
        assets_dir / "closing_backdrop.png",
        assets_dir / "closing_backdrop.jpg",
        Path("closing_backdrop.png"),
        Path("closing_backdrop.jpg"),
        DEFAULT_CLOSING_BACKDROP,
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _add_centered_line(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size,
    bold: bool,
    font_name: str,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=WHITE, name=font_name)


def add_closing_slide(
    prs: Presentation,
    report_date: str,
    *,
    assets_dir: Path,
    backdrop_path: Path | None = None,
    slide_number: int = 11,
) -> None:
    del report_date, slide_number
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])

    backdrop = resolve_closing_backdrop(assets_dir, backdrop_path)
    if backdrop is not None:
        picture = slide.shapes.add_picture(
            str(backdrop),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        _send_shape_to_back(picture)

    _add_centered_line(
        slide,
        text="We are imagining mobility with you",
        left=0.7,
        top=4.05,
        width=11.9,
        height=1.5,
        size=CLOSING_HEADLINE_SIZE,
        bold=True,
        font_name=FONT_MAJOR,
    )
    _add_centered_line(
        slide,
        text="Let's collaborate",
        left=1.2,
        top=5.45,
        width=10.9,
        height=0.55,
        size=CLOSING_SUBLINE_SIZE,
        bold=False,
        font_name=FONT_BODY,
    )
