#!/usr/bin/env python3
"""
Generate CES PFS CSAR (Non-STLA) Weekly Status Report (PowerPoint).

Uses templates/CES_CSAR_WSR_Template.pptx for brand fonts (Work Sans), colours, and
table styles extracted from the reference deck.

================================================================================
TODO — HARDCODED SECTIONS (not sourced from data.xlsm yet)
================================================================================
1. Presenter name on title slide
2. Agenda slide bullets
3. MOM & Action Items (meeting-specific narrative)
4. DCR status weekly narrative paragraphs (WK22–WK24) — fallback when graph Remarks empty
5. Risks & Mitigation content
6. Issues slide content
7. Quarterly planning chart values fallback
8. Closing slide backdrop image path (optional --closing-image)
================================================================================
SOURCED FROM data.xlsm (Non STLA, Visibility, DDP_Plan, planning, graph sheets)
================================================================================
- Pending eval/impl DCR tables
- DCR status summary callout boxes
- DDP MS4-5 slide
- Eval handoff slide (from Non_STLA Planning — Eval COmpleted column)
- Discussion points (at-risk DCRs)
- Charts and graph totals
================================================================================
"""

from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from wsr_charts import save_evaluation_chart, save_implementation_chart
from wsr_common import (
    DEFAULT_DATA_FILE,
    discussion_points,
    ddp_ms45_items,
    eval_handoff_items,
    get_evaluation_data,
    get_implementation_data,
    graph_week_capacity,
    load_ddp_plan,
    load_graph_summary,
    load_non_stla_planning,
    load_tracker,
    load_visibility,
    pending_items,
    pending_week_for_chart,
    summary_callouts,
    tracker_lookup,
    tracker_rows_lookup,
    week_remarks,
)
from wsr_style import (
    CLOSING_HEADLINE_SIZE,
    CLOSING_SUBLINE_SIZE,
    DEFAULT_CLOSING_BACKDROP,
    DEFAULT_TEMPLATE,
    FONT_BODY,
    FONT_MAJOR,
    FOOTER_DATE_HEIGHT,
    FOOTER_DATE_LEFT,
    FOOTER_DATE_TOP,
    FOOTER_DATE_WIDTH,
    FOOTER_NUMBER_HEIGHT,
    FOOTER_NUMBER_LEFT,
    FOOTER_NUMBER_TOP,
    FOOTER_NUMBER_WIDTH,
    FOOTER_SIZE,
    TEXT_DARK,
    TEXT_MUTED,
    TITLE_SIZE,
    WHITE,
    add_number_badge,
    delete_all_slides,
    find_placeholder,
    set_run_font,
    set_slide_footer,
    set_slide_title,
    style_title_date_run,
    style_agenda_run,
    style_body_run,
    style_table_cells,
    style_title_run,
)

LAYOUT_OPENING = 13
LAYOUT_CONTENT = 3

HARDCODED_PRESENTER = "Chitharenjan Nair"
HARDCODED_AGENDA = [
    "MOM & Action Items",
    "DCR Status",
    "Discussion Points",
    "Issues and Risks",
]
HARDCODED_MOM_ROWS = [
    {
        "discussion": "DCR 18922156 Phase 35 OBD REAL Regen Trigger Issue (Core 2 DCR). ",
        "action": "Implementation deferred to Q4'26.",
        "status": "Eval completed, Implementation deferred.",
        "owner": "PFS Team",
        "closure": "18-06-2026",
        "remarks": "Yue and Dhanraj confirmed that OBD REAL Regen Trigger issue for Core 2 program "
        "is not urgent in Q3'26. PFS Team need to work on the new Defect DCR for "
        "AftSootFiltFrqntRegenDiag component for both Eval and Impl.",
    }
]
AGENDA_BADGE_SIZE = 1.04
AGENDA_LAYOUT = [
    {"badge_top": 1.22, "text_left": 1.72, "text_top": 1.55},
    {"badge_top": 2.26, "text_left": 1.66, "text_top": 2.65},
    {"badge_top": 3.29, "text_left": 1.66, "text_top": 3.57},
    {"badge_top": 4.27, "text_left": 1.72, "text_top": 4.59},
]
HARDCODED_WEEKLY_NARRATIVE = {
    22: "WK:22 — 1 DCR: L2 Comments fixing",
    23: "WK:23 — Same as WK22 | 1 DCR: Awaiting CDS Review | 2 DCR's Awaiting L2 Review",
    24: "WK:24 — 12 DCR Closure pending | 5 DCR for L2 Reverification | 1 DCR WIP | 6 DCR Eval Closed",
}
FALLBACK_QUARTERLY_PLANNING = {
    "available_hours": 7514,
    "planned_pct": 90,
    "planned_hours": 6754,
    "resources": 17,
}
DEFAULT_PLANNING_BOOK = Path(__file__).parent / "Book2.xlsx"


def _as_float(value) -> float | None:
    try:
        if value in (None, ""):
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def load_quarterly_planning(planning_book: str | Path | None = None) -> dict[str, int]:
    workbook_path = Path(planning_book) if planning_book else DEFAULT_PLANNING_BOOK
    if not workbook_path.exists():
        return FALLBACK_QUARTERLY_PLANNING.copy()

    wb = load_workbook(workbook_path, data_only=True)
    ws = wb[wb.sheetnames[0]]

    available_hours = _as_float(ws["H34"].value)
    planned_hours = _as_float(ws["H36"].value)
    planned_pct_raw = _as_float(ws["H37"].value)
    resources_raw = _as_float(ws["I44"].value)

    if None in (available_hours, planned_hours, planned_pct_raw, resources_raw):
        return FALLBACK_QUARTERLY_PLANNING.copy()

    return {
        "available_hours": int(round(available_hours)),
        "planned_hours": int(round(planned_hours)),
        "planned_pct": int(round(planned_pct_raw * 100 if planned_pct_raw <= 1 else planned_pct_raw)),
        "resources": int(round(resources_raw)),
    }


def _new_content_slide(
    prs: Presentation,
    title: str,
    report_date: str,
    slide_number: int,
    *,
    title_size=None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    if title_size is not None:
        set_slide_title(slide, title, size=title_size)
    else:
        set_slide_title(slide, title)
    set_slide_footer(slide, report_date, slide_number)
    return slide


def _add_title_slide(prs: Presentation, report_date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_OPENING])

    title_ph = find_placeholder(slide, idx=0)
    if title_ph is not None:
        text_frame = title_ph.text_frame
        text_frame.clear()
        line1 = text_frame.paragraphs[0].add_run()
        line1.text = "CES PFS CSAR (Non- STLA) and CORE 2"
        style_title_run(line1)
        line2 = text_frame.add_paragraph().add_run()
        line2.text = "Weekly Status Report"
        set_run_font(line2, size=TITLE_SIZE, bold=False, color=TEXT_DARK, name=FONT_MAJOR)

    date_ph = find_placeholder(slide, idx=16)
    if date_ph is not None:
        date_ph.text = report_date
        for paragraph in date_ph.text_frame.paragraphs:
            for run in paragraph.runs:
                style_title_date_run(run)

    presenter_ph = find_placeholder(slide, idx=23)
    if presenter_ph is not None:
        presenter_ph.text = HARDCODED_PRESENTER
        for paragraph in presenter_ph.text_frame.paragraphs:
            for run in paragraph.runs:
                style_body_run(run)

    set_slide_footer(slide, report_date, 1)


def _add_agenda_slide(prs: Presentation, report_date: str):
    slide = _new_content_slide(prs, "Agenda", report_date, 2)
    for idx, (item, layout) in enumerate(zip(HARDCODED_AGENDA, AGENDA_LAYOUT), start=1):
        add_number_badge(
            slide,
            str(idx),
            left=0.67,
            top=layout["badge_top"],
            accent=idx >= 3,
            size=AGENDA_BADGE_SIZE,
        )
        box = slide.shapes.add_textbox(
            Inches(layout["text_left"]),
            Inches(layout["text_top"]),
            Inches(4.3),
            Inches(0.5),
        )
        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].add_run()
        run.text = item
        style_agenda_run(run)


def _add_mom_slide(prs: Presentation, report_date: str):
    slide = _new_content_slide(prs, "MOM of 11/06,18/06", report_date, 3)
    headers = [
        "Sr. No.",
        "Discussion points",
        "Action Items",
        "Status",
        "Ownership",
        "Action closure date",
        "Remarks",
    ]
    rows = [
        [
            str(idx),
            row["discussion"],
            row["action"],
            row["status"],
            row["owner"],
            row["closure"],
            row["remarks"],
        ]
        for idx, row in enumerate(HARDCODED_MOM_ROWS, start=1)
    ]
    _add_table(slide, headers, rows, top=1.0, col_widths=[0.55, 2.2, 1.5, 1.3, 0.85, 1.1, 2.35])


def _add_dcr_status_slide(
    prs: Presentation,
    report_date: str,
    impl_chart: Path,
    eval_chart: Path,
    summary: dict,
    callouts: dict,
    chart_week: int,
    pending_week: int,
    eval_data,
    impl_data,
):
    slide = _new_content_slide(
        prs,
        "DCR status Q3'26 – CSAR (Non-STLA)  & Core 2 program - PFS (Evaluation and Implementation)",
        report_date,
        4,
    )

    slide.shapes.add_picture(str(impl_chart), Inches(0.12), Inches(1.05), width=Inches(6.45))
    slide.shapes.add_picture(str(eval_chart), Inches(6.62), Inches(1.05), width=Inches(6.45))

    callout = slide.shapes.add_textbox(Inches(0.2), Inches(4.7), Inches(6.2), Inches(2.1))
    tf = callout.text_frame
    tf.word_wrap = True
    lines = [
        f"Total DCR's planned: {callouts['total_planned']}",
        f"{callouts['csar']} | {callouts['core2']} | "
        f"{callouts['ecm_testing']} | {callouts['ddp_testing']}",
        f"{callouts['eval_planned']} | {callouts['impl_planned']}",
        f"{callouts['rejected']} | {callouts['deferred']}",
        f"Data totals — Eval: {summary.get('eval_baseline')}/{summary.get('eval_revised')}/"
        f"{summary.get('eval_completed')} | Impl: {summary.get('impl_baseline')}/"
        f"{summary.get('impl_revised')}/{summary.get('impl_completed')}",
    ]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        set_run_font(run, size=Pt(8), color=TEXT_DARK, name=FONT_BODY)

    notes = slide.shapes.add_textbox(Inches(6.7), Inches(4.7), Inches(6.1), Inches(2.1))
    ntf = notes.text_frame
    ntf.word_wrap = True
    note_lines = []
    for week in (pending_week - 1, pending_week, chart_week):
        hardcoded = HARDCODED_WEEKLY_NARRATIVE.get(week)
        from_data_eval = week_remarks(eval_data, week)
        from_data_impl = week_remarks(impl_data, week)
        if from_data_eval or from_data_impl:
            note_lines.append(f"WK:{week} — Eval: {from_data_eval or '-'} | Impl: {from_data_impl or '-'}")
        elif hardcoded:
            note_lines.append(hardcoded)
    for i, line in enumerate(note_lines):
        para = ntf.paragraphs[0] if i == 0 else ntf.add_paragraph()
        run = para.add_run()
        run.text = line
        set_run_font(run, size=Pt(8), color=TEXT_MUTED, name=FONT_BODY)


def _add_pending_slide(
    prs: Presentation,
    title: str,
    report_date: str,
    slide_number: int,
    items: list[dict],
    mode: str,
):
    slide = _new_content_slide(prs, title, report_date, slide_number)

    if mode == "evaluation":
        headers = ["Sr No", "DCR ID", "Summary", "Current Status", "Eval Closure date", "Support Required"]
        rows = [
            [
                str(i + 1),
                str(item["dcr_id"]),
                item["summary"],
                item["status"],
                item["closure_date"],
                item["support"],
            ]
            for i, item in enumerate(items)
        ]
        widths = [0.5, 0.75, 2.7, 1.6, 1.0, 2.6]
    else:
        headers = ["Sr No", "DCR ID", "Summary", "Current Status", "Remarks"]
        rows = [
            [
                str(i + 1),
                str(item["dcr_id"]),
                item["summary"],
                item["status"],
                item["remarks"],
            ]
            for i, item in enumerate(items)
        ]
        widths = [0.5, 0.75, 3.2, 1.8, 2.9]

    if not rows:
        rows = [["-", "-", "No open items found with current filters", "-", "-"]]
        if mode == "evaluation":
            rows[0].append("-")

    _add_table(slide, headers, rows, top=0.95, col_widths=widths)


def _add_ddp_slide(prs: Presentation, report_date: str, items: list[dict]):
    slide = _new_content_slide(prs, "PFS DDP Details MS 4-5", report_date, 7)
    headers = [
        "Sr.No",
        "DCR ID",
        "DCR Summary",
        "Plan Dates",
        "Appeared Dates",
        "Program",
        "Dependencies",
        "Remarks",
    ]
    rows = [
        [
            str(item["sr_no"]),
            item["dcr_id"],
            item["summary"],
            item["plan_date"],
            item["appeared_date"],
            item["program"],
            item.get("dependencies", "-"),
            item["remarks"],
        ]
        for item in items
    ]
    if not rows:
        rows = [["-", "-", "No DDP MS4-5 rows matched filters", "-", "-", "-", "-", "-"]]
    _add_table(
        slide,
        headers,
        rows,
        top=0.95,
        col_widths=[0.45, 0.7, 2.0, 0.9, 0.9, 0.65, 1.55, 1.9],
    )


def _add_handoff_slide(prs: Presentation, report_date: str, items: list[dict]):
    slide = _new_content_slide(prs, "Q3-2026 – Eval Handoff from onsite", report_date, 8)
    headers = ["Sr. No.", "DCR ID", "Summary", "Evaluator", "Eval Handoff Date", "Remark"]
    rows = []
    for idx, item in enumerate(items, start=1):
        rows.append(
            [
                str(idx),
                str(item["dcr_id"]),
                str(item.get("summary", "-")),
                item.get("evaluator", "-"),
                item.get("handoff_date", "-"),
                item.get("remark", "-"),
            ]
        )
    if not rows:
        rows = [["-", "-", "No eval handoff rows for this month", "-", "-", "-"]]
    _add_table(slide, headers, rows, top=0.95, col_widths=[0.55, 0.75, 2.9, 0.9, 1.2, 1.9])


def _add_discussion_slide(prs: Presentation, report_date: str, items: list[dict]):
    slide = _new_content_slide(prs, "Discussion Points", report_date, 9)
    headers = ["#", "Descriptions", "Program", "Priority", "Plan completion dates", "Remarks"]
    rows = [
        [
            str(i + 1),
            item["description"],
            item["program"],
            str(item["priority"]),
            item["plan_date"],
            item["remarks"],
        ]
        for i, item in enumerate(items)
    ]
    if not rows:
        rows = [["-", "No high-priority discussion points found", "-", "-", "-", "-"]]
    _add_table(slide, headers, rows, top=0.95, col_widths=[0.35, 3.0, 0.75, 0.7, 1.2, 2.5])


def _add_impact_legend(slide):
    """Risk/issues impact legend (matches reference slide 10 footer)."""
    legend_items = [
        (RGBColor(0xFF, 0x00, 0x00), "High Impact / High Possibility", 0.76, 0.98, 6.38),
        (RGBColor(0xFF, 0xC0, 0x00), "Medium Impact / Medium Possibility", 3.32, 3.48, 6.38),
        (RGBColor(0x92, 0xD0, 0x50), "Low Impact / Low Possibility", 6.23, 6.45, 6.39),
    ]
    for color, label, square_left, text_left, square_top in legend_items:
        square = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(square_left),
            Inches(square_top),
            Inches(0.21),
            Inches(0.21),
        )
        square.fill.solid()
        square.fill.fore_color.rgb = color
        square.line.fill.background()

        box = slide.shapes.add_textbox(Inches(text_left), Inches(6.35), Inches(2.6), Inches(0.29))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = label
        style_body_run(run)


def _add_risks_slide(prs: Presentation, report_date: str):
    slide = _new_content_slide(prs, "Risks & Mitigation Plan", report_date, 10)
    headers = ["#", "Risk", "Impact", "Risk Mitigation / Contingency"]
    rows = [["", "NA", "", ""]]
    _add_table(slide, headers, rows, top=0.85, col_widths=[0.5, 3.5, 1.2, 3.5])

    issue_headers = ["#", "Issues", "Impact", "Contingency action"]
    issue_rows = [["-", "-", "-", "-"]]
    _add_table(slide, issue_headers, issue_rows, top=5.55, col_widths=[0.5, 3.5, 1.2, 3.5])

    _add_impact_legend(slide)


def _add_planning_slide(prs: Presentation, report_date: str, qp: dict[str, int]):
    slide = _new_content_slide(prs, "Quarterly Planning 2026-Non STLA", report_date, 11)

    chart_data = CategoryChartData()
    chart_data.categories = ["Q3 Actual Available", f"{qp['planned_pct']}% of Q3 is Planned"]
    chart_data.add_series("", (qp["available_hours"], qp["planned_hours"]))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.29),
        Inches(0.93),
        Inches(11.20),
        Inches(4.91),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_title = True
    title_frame = chart.chart_title.text_frame
    title_frame.text = "PFS Quarterly Planning 2026"
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            set_run_font(run, size=Pt(14), bold=True, color=TEXT_DARK, name=FONT_BODY)

    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x92, 0xD0, 0x50)
    series.data_labels.number_format = "0"

    bandwidth = slide.shapes.add_textbox(Inches(0.59), Inches(5.85), Inches(10.85), Inches(0.40))
    bw_run = bandwidth.text_frame.paragraphs[0].add_run()
    bw_run.text = f"**Bandwidth is planned with {qp['resources']} Resources"
    style_body_run(bw_run, bold=True)

    note = slide.shapes.add_textbox(Inches(0.59), Inches(6.25), Inches(12.18), Inches(0.40))
    note_run = note.text_frame.paragraphs[0].add_run()
    note_run.text = (
        "Note: The Estimations for Q3 Planning are in progress, "
        "these are high level tentative estimations"
    )
    style_body_run(note_run)


def _send_shape_to_back(shape) -> None:
    element = shape._element
    parent = element.getparent()
    parent.remove(element)
    parent.insert(0, element)


def _clear_placeholder_text(slide, idx: int) -> None:
    placeholder = find_placeholder(slide, idx=idx)
    if placeholder is not None:
        placeholder.text = ""


def _resolve_closing_backdrop(assets_dir: Path, override: Path | None = None) -> Path | None:
    if override is not None and override.exists():
        return override
    candidates = [
        assets_dir / "closing_backdrop.png",
        assets_dir / "closing_backdrop.jpg",
        Path("closing_backdrop.png"),
        Path("closing_backdrop.jpg"),
        DEFAULT_CLOSING_BACKDROP,
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _set_closing_footer(slide, report_date: str, slide_number: int) -> None:
    for left, top, width, height, text in (
        (FOOTER_DATE_LEFT, FOOTER_DATE_TOP, FOOTER_DATE_WIDTH, FOOTER_DATE_HEIGHT, report_date),
        (FOOTER_NUMBER_LEFT, FOOTER_NUMBER_TOP, FOOTER_NUMBER_WIDTH, FOOTER_NUMBER_HEIGHT, str(slide_number)),
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = box.text_frame
        text_frame.clear()
        text_frame.margin_left = text_frame.margin_right = text_frame.margin_top = text_frame.margin_bottom = 0
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        run = paragraph.add_run()
        run.text = text
        set_run_font(run, size=FOOTER_SIZE, color=WHITE, name=FONT_BODY)


def _add_closing_slide(
    prs: Presentation,
    report_date: str,
    *,
    assets_dir: Path,
    backdrop_path: Path | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_OPENING])

    backdrop = _resolve_closing_backdrop(assets_dir, backdrop_path)
    if backdrop is not None:
        picture = slide.shapes.add_picture(
            str(backdrop),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        _send_shape_to_back(picture)

    for placeholder_idx in (0, 16, 23):
        _clear_placeholder_text(slide, placeholder_idx)

    headline_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.05), Inches(11.9), Inches(1.5))
    headline_frame = headline_box.text_frame
    headline_frame.clear()
    headline_frame.word_wrap = True
    headline_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    headline_paragraph = headline_frame.paragraphs[0]
    headline_paragraph.alignment = PP_ALIGN.CENTER
    headline_run = headline_paragraph.add_run()
    headline_run.text = "We are imagining mobility with you"
    set_run_font(
        headline_run,
        size=CLOSING_HEADLINE_SIZE,
        bold=True,
        color=WHITE,
        name=FONT_MAJOR,
    )

    subline_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.45), Inches(10.9), Inches(0.55))
    subline_frame = subline_box.text_frame
    subline_frame.clear()
    subline_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    subline_paragraph = subline_frame.paragraphs[0]
    subline_paragraph.alignment = PP_ALIGN.CENTER
    subline_run = subline_paragraph.add_run()
    subline_run.text = "Lets collaborate"
    set_run_font(
        subline_run,
        size=CLOSING_SUBLINE_SIZE,
        bold=False,
        color=WHITE,
        name=FONT_BODY,
    )

    _set_closing_footer(slide, report_date, 12)


def _add_table(slide, headers, rows, top=1.0, col_widths=None):
    row_height = min(0.34 * (len(rows) + 1), 5.8)
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(0.25),
        Inches(top),
        Inches(12.85),
        Inches(row_height),
    )
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)

    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)

    style_table_cells(table)


def generate_report(
    output_path: str | Path = "WSR_Report.pptx",
    data_file: str = DEFAULT_DATA_FILE,
    chart_week: int = 25,
    report_date: str | None = None,
    assets_dir: str | Path = "report_assets",
    template_path: str | Path = DEFAULT_TEMPLATE,
    closing_image: str | Path | None = None,
    planning_book: str | Path | None = None,
) -> Path:
    assets_dir = Path(assets_dir)
    assets_dir.mkdir(exist_ok=True)
    closing_image_path = Path(closing_image) if closing_image else None
    planning_book_path = Path(planning_book) if planning_book else None
    pending_week = pending_week_for_chart(chart_week)

    if report_date is None:
        report_date = datetime.now().strftime("%d-%m-%Y")

    impl_chart = save_implementation_chart(assets_dir / "implementation_chart.png", data_file=data_file)
    eval_chart = save_evaluation_chart(assets_dir / "evaluation_chart.png", data_file=data_file)

    summary = load_graph_summary(data_file)
    callouts = summary_callouts(data_file)
    eval_data = get_evaluation_data(data_file=data_file)
    impl_data = get_implementation_data(data_file=data_file)
    tracker = load_tracker(data_file)
    visibility = load_visibility(data_file)
    ddp = load_ddp_plan(data_file)
    planning = load_non_stla_planning(data_file)
    tracker_map = tracker_lookup(tracker)
    tracker_rows = tracker_rows_lookup(tracker)

    eval_limit = graph_week_capacity(data_file, pending_week, "evaluation")
    eval_pending = pending_items(
        visibility, tracker_rows, mode="evaluation", pending_week=pending_week, limit=eval_limit
    )
    impl_pending = pending_items(
        visibility, tracker_rows, mode="implementation", pending_week=pending_week, limit=2
    )
    ddp_items = ddp_ms45_items(ddp, tracker_map)
    handoff_items = eval_handoff_items(planning, tracker_map, report_date)
    discussion = discussion_points(visibility, tracker_map)
    quarterly_planning = load_quarterly_planning(planning_book_path)

    prs = Presentation(str(template_path))
    delete_all_slides(prs)

    _add_title_slide(prs, report_date)
    _add_agenda_slide(prs, report_date)
    _add_mom_slide(prs, report_date)
    _add_dcr_status_slide(
        prs,
        report_date,
        impl_chart,
        eval_chart,
        summary,
        callouts,
        chart_week,
        pending_week,
        eval_data,
        impl_data,
    )
    _add_pending_slide(
        prs,
        f"Q3-2026 – Evaluations pending for closure for week {pending_week}",
        report_date,
        5,
        eval_pending,
        mode="evaluation",
    )
    _add_pending_slide(
        prs,
        f"Q3-2026 – Implementation pending for closure for week {pending_week}",
        report_date,
        6,
        impl_pending,
        mode="implementation",
    )
    _add_ddp_slide(prs, report_date, ddp_items)
    _add_handoff_slide(prs, report_date, handoff_items)
    _add_discussion_slide(prs, report_date, discussion)
    _add_risks_slide(prs, report_date)
    _add_planning_slide(prs, report_date, quarterly_planning)
    _add_closing_slide(
        prs,
        report_date,
        assets_dir=assets_dir,
        backdrop_path=closing_image_path,
    )

    output_path = Path(output_path)
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate WSR PowerPoint report from data.xlsm")
    parser.add_argument("--data", default=DEFAULT_DATA_FILE, help="Path to Excel workbook")
    parser.add_argument("--output", default="WSR_Report.pptx", help="Output PowerPoint path")
    parser.add_argument(
        "--week",
        type=int,
        default=25,
        help="Chart week number (pending tables use week-1, matching reference PDF)",
    )
    parser.add_argument("--date", default=None, help="Report date label (dd-mm-yyyy)")
    parser.add_argument("--assets-dir", default="report_assets", help="Directory for chart images")
    parser.add_argument(
        "--closing-image",
        default=None,
        help="Closing slide backdrop image (default: report_assets/closing_backdrop.png)",
    )
    parser.add_argument(
        "--planning-book",
        default=None,
        help="Quarterly planning workbook for slide 11 (default: ./Book2.xlsx)",
    )
    parser.add_argument(
        "--template",
        default=str(DEFAULT_TEMPLATE),
        help="Branded PowerPoint template (default: templates/CES_CSAR_WSR_Template.pptx)",
    )
    args = parser.parse_args()

    output = generate_report(
        output_path=args.output,
        data_file=args.data,
        chart_week=args.week,
        report_date=args.date,
        assets_dir=args.assets_dir,
        template_path=args.template,
        closing_image=args.closing_image,
        planning_book=args.planning_book,
    )
    print(f"Report generated: {output}")
    print(f"Template: {args.template}")
    print(f"Chart week: {args.week} | Pending tables week: {pending_week_for_chart(args.week)}")


if __name__ == "__main__":
    main()
