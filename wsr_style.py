"""Visual theme for WSR PowerPoint reports (CES CSAR Cummins template)."""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from wsr.resources import resource_path

DEFAULT_TEMPLATE = resource_path("templates", "CES_CSAR_WSR_Template.pptx")

# Cummins / CES deck palette (from template theme1.xml)
ACCENT_LIME = RGBColor(0xB0, 0xFF, 0x45)      # accent1
TEXT_DARK = RGBColor(0x16, 0x17, 0x18)       # dk1
TEXT_MUTED = RGBColor(0x4D, 0x51, 0x54)      # dk2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUMMARY_TABLE_FILL = RGBColor(0xC6, 0xE0, 0xB4)

FONT_MAJOR = "Work Sans Medium"
FONT_BODY = "Work Sans"

TABLE_STYLE_ID = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"

SLIDE_WIDTH_IN = 13.333
TABLE_MARGIN_IN = 0.4
TABLE_WIDTH_IN = SLIDE_WIDTH_IN - (2 * TABLE_MARGIN_IN)
TABLE_LEFT_IN = TABLE_MARGIN_IN
TABLE_HEADER_ROW_HEIGHT_IN = 0.46
TABLE_TOP_GAP_IN = 0.22
TABLE_TOP_MIN_IN = 1.45

TITLE_SIZE = Pt(32)
TITLE_DATE_SIZE = Pt(14)
SECTION_TITLE_SIZE = Pt(24)
BODY_SIZE = Pt(12)
AGENDA_ITEM_SIZE = Pt(24)
TABLE_HEADER_SIZE = Pt(10.5)
TABLE_BODY_SIZE = Pt(10)
FOOTER_SIZE = Pt(9)
CLOSING_HEADLINE_SIZE = Pt(44)
CLOSING_SUBLINE_SIZE = Pt(20)
DEFAULT_CLOSING_BACKDROP = resource_path("assets", "closing_backdrop.png")

# Bottom-right footer positions from CES_CSAR reference deck (date, then slide number).
FOOTER_DATE_LEFT = 11.17
FOOTER_DATE_TOP = 6.92
FOOTER_DATE_WIDTH = 1.06
FOOTER_DATE_HEIGHT = 0.27
FOOTER_NUMBER_LEFT = 12.24
FOOTER_NUMBER_TOP = 6.92
FOOTER_NUMBER_WIDTH = 0.51
FOOTER_NUMBER_HEIGHT = 0.26
FOOTER_MIN_TOP = Inches(6.5)


def set_run_font(
    run,
    *,
    size=BODY_SIZE,
    bold=False,
    color=TEXT_DARK,
    name=FONT_BODY,
):
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def style_title_run(run):
    set_run_font(run, size=TITLE_SIZE, bold=True, color=TEXT_DARK, name=FONT_MAJOR)


def style_section_title_run(run):
    set_run_font(run, size=SECTION_TITLE_SIZE, bold=True, color=TEXT_DARK, name=FONT_MAJOR)


def style_body_run(run, *, bold=False, color=TEXT_DARK):
    set_run_font(run, size=BODY_SIZE, bold=bold, color=color, name=FONT_BODY)


def style_title_date_run(run):
    set_run_font(run, size=TITLE_DATE_SIZE, bold=False, color=TEXT_DARK, name=FONT_BODY)


def style_agenda_run(run):
    set_run_font(run, size=AGENDA_ITEM_SIZE, color=TEXT_DARK, name=FONT_BODY)


def style_footer_run(run):
    set_run_font(run, size=FOOTER_SIZE, color=TEXT_MUTED, name=FONT_BODY)


def delete_all_slides(prs) -> None:
    for index in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[index]
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[index]


def find_placeholder(slide, idx: int | None = None, name_contains: str | None = None):
    for shape in slide.placeholders:
        if idx is not None and shape.placeholder_format.idx == idx:
            return shape
        if name_contains and name_contains.lower() in shape.name.lower():
            return shape
    return None


def _footer_placeholder(slide, idx_candidates: list[int], name_hint: str):
    for idx in idx_candidates:
        ph = find_placeholder(slide, idx=idx)
        if ph is not None and ph.top >= FOOTER_MIN_TOP:
            return ph
    ph = find_placeholder(slide, name_contains=name_hint)
    if ph is not None and ph.top >= FOOTER_MIN_TOP:
        return ph
    return None


def _write_footer_text(shape, text: str) -> None:
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            style_footer_run(run)


def _add_footer_textbox(slide, left: float, top: float, width: float, height: float, text: str):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    style_footer_run(run)


def set_slide_footer(slide, report_date: str, slide_number: int | None = None):
    date_ph = _footer_placeholder(slide, [13, 17, 16, 14], "Date")
    if date_ph is not None:
        _write_footer_text(date_ph, report_date)
    else:
        _add_footer_textbox(
            slide,
            FOOTER_DATE_LEFT,
            FOOTER_DATE_TOP,
            FOOTER_DATE_WIDTH,
            FOOTER_DATE_HEIGHT,
            report_date,
        )

    if slide_number is not None:
        num_ph = _footer_placeholder(slide, [10, 19, 22, 12], "Slide Number")
        if num_ph is not None:
            _write_footer_text(num_ph, str(slide_number))
        else:
            _add_footer_textbox(
                slide,
                FOOTER_NUMBER_LEFT,
                FOOTER_NUMBER_TOP,
                FOOTER_NUMBER_WIDTH,
                FOOTER_NUMBER_HEIGHT,
                str(slide_number),
            )


def set_slide_title(slide, title: str, *, size=SECTION_TITLE_SIZE):
    title_ph = find_placeholder(slide, idx=0) or find_placeholder(slide, name_contains="Title")
    if title_ph is None:
        return
    title_ph.text = title
    for paragraph in title_ph.text_frame.paragraphs:
        for run in paragraph.runs:
            set_run_font(run, size=size, bold=True, color=TEXT_DARK, name=FONT_MAJOR)


def content_top_below_title(
    slide,
    *,
    gap: float = TABLE_TOP_GAP_IN,
    minimum: float = TABLE_TOP_MIN_IN,
) -> float:
    title_ph = find_placeholder(slide, idx=0) or find_placeholder(slide, name_contains="Title")
    if title_ph is not None and title_ph.width.inches > 1.0:
        bottom = title_ph.top.inches + title_ph.height.inches
        return max(minimum, bottom + gap)
    return minimum


def raise_slide_title(slide, top_in: float = 0.42) -> None:
    title_ph = find_placeholder(slide, idx=0) or find_placeholder(slide, name_contains="Title")
    if title_ph is not None:
        left = title_ph.left
        width = title_ph.width
        height = title_ph.height
        title_ph.top = Inches(top_in)
        title_ph.left = left
        title_ph.width = width
        title_ph.height = height


def style_key_value_table(table) -> None:
    apply_table_style(table)
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SUMMARY_TABLE_FILL
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            is_label = col_idx == 0
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                if not paragraph.runs:
                    paragraph.text = paragraph.text or ""
                for run in paragraph.runs:
                    set_run_font(
                        run,
                        size=Pt(9) if is_label else Pt(9.5),
                        bold=is_label,
                        color=TEXT_DARK,
                        name=FONT_BODY,
                    )
                    if is_label:
                        run.font.bold = True


def min_header_column_width(header: str) -> float:
    words = [word for word in str(header).split() if word]
    if not words:
        return 0.5
    longest_word = max(len(word) for word in words)
    return max(0.52, longest_word * 0.095)


def fit_table_column_widths(headers: list[str], relative_widths: list[float] | None = None) -> list[float]:
    if not headers:
        return []
    if not relative_widths:
        relative_widths = [1.0] * len(headers)
    mins = [min_header_column_width(header) for header in headers]
    widths = [max(relative, minimum) for relative, minimum in zip(relative_widths, mins)]
    total = sum(widths)
    if total <= 0:
        even = TABLE_WIDTH_IN / len(headers)
        return [even] * len(headers)
    scale = TABLE_WIDTH_IN / total
    return [width * scale for width in widths]


def apply_table_style(table) -> None:
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    for child in list(tbl_pr):
        if child.tag == qn("a:tableStyleId"):
            tbl_pr.remove(child)
    style = OxmlElement("a:tableStyleId")
    style.text = TABLE_STYLE_ID
    tbl_pr.append(style)
    if tbl_pr.get("firstRow") is None:
        tbl_pr.set("firstRow", "1")
    if tbl_pr.get("bandRow") is None:
        tbl_pr.set("bandRow", "1")


def style_table_cells(table, *, header_rows: int = 1):
    apply_table_style(table)
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

            is_header = row_idx < header_rows
            text_frame = cell.text_frame
            text_frame.word_wrap = True
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
                if not paragraph.runs:
                    paragraph.text = paragraph.text or ""
                for run in paragraph.runs:
                    set_run_font(
                        run,
                        size=TABLE_HEADER_SIZE if is_header else TABLE_BODY_SIZE,
                        bold=is_header,
                        color=TEXT_DARK,
                        name=FONT_BODY,
                    )


def add_number_badge(
    slide,
    number: str,
    left: float,
    top: float,
    accent: bool = False,
    *,
    size: float = 0.42,
):
    dimension = Inches(size)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), dimension, dimension)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIME if accent else TEXT_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number
    badge_font = Pt(24) if size >= 1.0 else Pt(14)
    set_run_font(
        run,
        size=badge_font,
        bold=True,
        color=WHITE if not accent else TEXT_DARK,
        name=FONT_MAJOR,
    )
    return shape
